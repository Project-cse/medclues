"""
MediChain+ — Professional PowerPoint Presentation Generator
Run this script to create the full PPT: python generate_ppt.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x1D, 0x37)   # Deep navy blue  (background)
CYAN        = RGBColor(0x00, 0xC2, 0xCB)   # Bright cyan     (accent)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # White           (text)
LIGHT_BLUE  = RGBColor(0xE8, 0xF4, 0xFF)   # Light blue      (subtle bg)
DARK_GRAY   = RGBColor(0x1E, 0x2D, 0x40)   # Dark gray       (card bg)
SOFT_WHITE  = RGBColor(0xCC, 0xE5, 0xFF)   # Soft blue-white (sub-text)
GREEN       = RGBColor(0x00, 0xE0, 0x96)   # Mint green      (check marks)
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)   # Orange          (highlight)

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # Blank layout


def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_multiline(slide, lines, left, top, width, height,
                  font_name="Calibri", font_size=16, bold=False,
                  color=WHITE, align=PP_ALIGN.LEFT, line_spacing=None):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.color.rgb = color
        if line_spacing:
            p.space_after = Pt(line_spacing)


def bg(slide, color=NAVY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=color)


def accent_bar(slide, y=6.8, color=CYAN):
    add_rect(slide, 0, y, 13.33, 0.08, fill_color=color)


def slide_number(slide, num, total=14, color=SOFT_WHITE):
    add_text(slide, f"{num} / {total}", 12.5, 7.1, 0.8, 0.3,
             font_size=9, color=color, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE SLIDE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)

# Left accent stripe
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)

# Glow circle decoration (top right)
add_rect(s, 10.5, -1.5, 4.5, 4.5, fill_color=DARK_GRAY)

# Tag line box
add_rect(s, 1.2, 1.0, 4.0, 0.45, fill_color=CYAN)
add_text(s, "  HEALTHCARE TECHNOLOGY  ", 1.2, 1.0, 4.0, 0.45,
         font_size=10, bold=True, color=NAVY, align=PP_ALIGN.CENTER)

# Main Title
add_text(s, "MediChain+", 1.2, 1.6, 9.0, 1.6,
         font_name="Calibri", font_size=72, bold=True, color=CYAN,
         align=PP_ALIGN.LEFT)

# Subtitle
add_text(s, "India's Smartest Hospital Network Platform", 1.2, 3.1, 9.5, 0.7,
         font_size=26, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

# Description
add_text(s, "Connecting Patients · Hospitals · Labs · Blood Banks · Emergency Services",
         1.2, 3.75, 10.0, 0.5, font_size=14, italic=True, color=SOFT_WHITE,
         align=PP_ALIGN.LEFT)

# Divider
add_rect(s, 1.2, 4.3, 5.0, 0.04, fill_color=CYAN)

# Bottom info
add_multiline(s, [
    "Presented by:  MediChain+ Development Team",
    "Platform:  Web Application  |  Stack: React · FastAPI · MongoDB",
    "April 2026"
], 1.2, 4.5, 8.0, 1.5, font_size=13, color=SOFT_WHITE)

accent_bar(s)
slide_number(s, 1)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=ORANGE)

add_text(s, "THE PROBLEM", 0.5, 0.3, 5.0, 0.55,
         font_size=11, bold=True, color=ORANGE)
add_text(s, "What We Are Solving", 0.5, 0.7, 10.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 2.5, 0.05, fill_color=ORANGE)

problems = [
    ("🏥", "No Central Platform",
     "Patients have no single place to find hospitals,\nlabs, blood banks, and doctors together."),
    ("🚑", "Emergency Confusion",
     "During emergencies, people waste critical\nminutes calling hospitals one by one."),
    ("🩸", "Blood Bank Mystery",
     "Patients and families have no way to\ncheck which blood bank has their needed blood group."),
    ("⏳", "Appointment Chaos",
     "Long queues, no online booking, no\ndigital token system — patients wait for hours."),
]

cols = [(0.5, 2.0), (3.55, 2.0), (6.6, 2.0), (9.65, 2.0)]
for i, (icon, title, desc) in enumerate(problems):
    x, y = cols[i]
    add_rect(s, x, y, 2.9, 4.8, fill_color=DARK_GRAY, line_color=ORANGE, line_width=0.5)
    add_text(s, icon, x + 0.1, y + 0.3, 2.7, 0.7,
             font_size=36, align=PP_ALIGN.CENTER)
    add_text(s, title, x + 0.1, y + 1.1, 2.7, 0.55,
             font_size=14, bold=True, color=ORANGE, align=PP_ALIGN.CENTER)
    add_text(s, desc, x + 0.15, y + 1.7, 2.6, 2.8,
             font_size=11, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

accent_bar(s, color=ORANGE)
slide_number(s, 2)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — THE SOLUTION
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)

add_text(s, "THE SOLUTION", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "Introducing MediChain+", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 3.5, 0.05, fill_color=CYAN)

add_text(s,
         "MediChain+ is a comprehensive, all-in-one healthcare platform that seamlessly connects "
         "patients with hospitals, doctors, diagnostic laboratories, blood banks, and emergency "
         "services — all in one powerful application.",
         0.5, 1.7, 12.3, 1.0, font_size=15, color=SOFT_WHITE)

# Feature pills
pills = [
    ("🏥 Hospital Discovery", CYAN),
    ("👨‍⚕️ Doctor Booking", GREEN),
    ("🔬 Lab Management", ORANGE),
    ("🩸 Blood Bank", RGBColor(0xFF,0x4D,0x6D)),
    ("🚨 Emergency Center", RGBColor(0xFF,0x6B,0x00)),
    ("🤖 AI Assistant", RGBColor(0x7B,0x61,0xFF)),
]

pill_positions = [
    (0.5, 2.95), (3.5, 2.95), (6.5, 2.95),
    (0.5, 4.05), (3.5, 4.05), (6.5, 4.05)
]
for (label, color), (px, py) in zip(pills, pill_positions):
    add_rect(s, px, py, 2.7, 0.75, fill_color=DARK_GRAY, line_color=color, line_width=1)
    add_text(s, label, px + 0.1, py + 0.1, 2.5, 0.55,
             font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)

# Right side statement
add_rect(s, 9.8, 2.7, 3.2, 3.9, fill_color=DARK_GRAY, line_color=CYAN, line_width=1)
add_text(s, "One Platform.\nEvery Need.\nEvery Patient.", 9.9, 2.9, 3.0, 2.0,
         font_size=22, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, "No more switching between\nmultiple apps. Everything\nyou need — in MediChain+.",
         9.9, 4.85, 3.0, 1.6, font_size=12, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

accent_bar(s)
slide_number(s, 3)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — KEY FEATURES (PART 1)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)

add_text(s, "CORE FEATURES", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "What MediChain+ Offers", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 3.5, 0.05, fill_color=CYAN)

features = [
    ("🏥", "Smart Hospital Discovery",
     "Locate real hospitals near you using live GPS.\nFilter by specialty, type, and distance.\nView, call, or navigate directly from the app.",
     CYAN),
    ("👨‍⚕️", "Doctor Appointment Booking",
     "Browse doctors by specialization and hospital.\nCheck real-time availability.\nBook appointments in under 60 seconds.",
     GREEN),
    ("🔬", "Laboratory & Test Booking",
     "Discover nearby diagnostic labs.\nFilter by test type and rating.\nBook your test slot instantly.",
     ORANGE),
    ("🩸", "Blood Bank Network",
     "View partnered blood banks with stock status.\nFilter by blood group availability.\nContact blood banks directly from the platform.",
     RGBColor(0xFF,0x4D,0x6D)),
]

rows = [(0.4, 2.0), (3.5, 2.0), (6.6, 2.0), (9.7, 2.0)]
for (icon, title, desc, color), (rx, ry) in zip(features, rows):
    add_rect(s, rx, ry, 3.0, 5.0, fill_color=DARK_GRAY, line_color=color, line_width=1)
    add_rect(s, rx, ry, 3.0, 0.12, fill_color=color)
    add_text(s, icon, rx + 0.1, ry + 0.25, 2.8, 0.7, font_size=34, align=PP_ALIGN.CENTER)
    add_text(s, title, rx + 0.1, ry + 1.05, 2.8, 0.7,
             font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, desc, rx + 0.15, ry + 1.8, 2.7, 3.0,
             font_size=11, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

accent_bar(s)
slide_number(s, 4)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — KEY FEATURES (PART 2)
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=GREEN)

add_text(s, "ADVANCED FEATURES", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=GREEN)
add_text(s, "What Sets Us Apart", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 3.5, 0.05, fill_color=GREEN)

adv = [
    ("🚨", "Emergency Center",
     "One-tap call to 108 Ambulance.\nAuto-sends your GPS location to emergency contacts.\nRoutes to the nearest available hospital.",
     RGBColor(0xFF,0x6B,0x00)),
    ("🤖", "AI Health Assistant",
     "Conversational AI trained on medical knowledge.\nAnswers health queries 24/7.\nProvides safe, evidence-based guidance.",
     RGBColor(0x7B,0x61,0xFF)),
    ("🛡️", "Multi-Role Admin System",
     "Super Admin manages all hospitals globally.\nDean Portal controls hospital-level operations.\nDoctor Portal for scheduling and patient management.",
     CYAN),
    ("📍", "Real-Time GPS Discovery",
     "Uses OpenStreetMap — not just our partner network.\nShows ALL nearby medical facilities.\nWorks even for non-partnered hospitals.",
     GREEN),
]

rows = [(0.4, 2.0), (3.5, 2.0), (6.6, 2.0), (9.7, 2.0)]
for (icon, title, desc, color), (rx, ry) in zip(adv, rows):
    add_rect(s, rx, ry, 3.0, 5.0, fill_color=DARK_GRAY, line_color=color, line_width=1)
    add_rect(s, rx, ry, 3.0, 0.12, fill_color=color)
    add_text(s, icon, rx + 0.1, ry + 0.25, 2.8, 0.7, font_size=34, align=PP_ALIGN.CENTER)
    add_text(s, title, rx + 0.1, ry + 1.05, 2.8, 0.7,
             font_size=13, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, desc, rx + 0.15, ry + 1.8, 2.7, 3.0,
             font_size=11, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

accent_bar(s, color=GREEN)
slide_number(s, 5)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — HOW IT WORKS
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)

add_text(s, "HOW IT WORKS", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "Simple. Fast. Reliable.", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 3.0, 0.05, fill_color=CYAN)

steps = [
    ("01", "Register / Login",
     "Patient creates a secure account in under 2 minutes with email verification."),
    ("02", "Find What You Need",
     "Search for hospitals, doctors, labs, or blood banks — filtered by your location and needs."),
    ("03", "Book in One Click",
     "Select your preferred doctor or lab, choose a time slot, and confirm your booking instantly."),
    ("04", "Get Your Token",
     "Receive a digital token number with live queue updates so you know exactly when it's your turn."),
    ("05", "Track & Manage",
     "View your appointment history, upcoming visits, and health records all in one dashboard."),
]

for i, (num, title, desc) in enumerate(steps):
    y = 1.9 + i * 1.0
    # number circle bg
    add_rect(s, 0.5, y + 0.05, 0.65, 0.65, fill_color=CYAN)
    add_text(s, num, 0.5, y + 0.05, 0.65, 0.65,
             font_size=16, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, title, 1.3, y + 0.05, 3.5, 0.45,
             font_size=14, bold=True, color=CYAN)
    add_text(s, desc, 1.3, y + 0.45, 8.5, 0.45,
             font_size=12, color=SOFT_WHITE)
    if i < 4:
        add_rect(s, 0.75, y + 0.72, 0.15, 0.35, fill_color=DARK_GRAY)

accent_bar(s)
slide_number(s, 6)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — USER ROLES
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=RGBColor(0x7B,0x61,0xFF))

add_text(s, "USER ROLES", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True,
         color=RGBColor(0x7B,0x61,0xFF))
add_text(s, "Designed for Everyone in Healthcare", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 4.5, 0.05, fill_color=RGBColor(0x7B,0x61,0xFF))

roles = [
    ("👤", "Patient",
     ["✦ Search nearby hospitals & labs",
      "✦ Book doctor appointments",
      "✦ Track appointment history",
      "✦ Access emergency services",
      "✦ View blood bank availability"],
     CYAN),
    ("🩺", "Doctor",
     ["✦ Manage daily appointment schedule",
      "✦ View patient booking details",
      "✦ Update availability status",
      "✦ Access patient appointment history",
      "✦ Secure login credentials via email"],
     GREEN),
    ("🏥", "Dean (Hospital Admin)",
     ["✦ Manage hospital doctors & staff",
      "✦ Add / edit / deactivate doctors",
      "✦ Monitor hospital-level bookings",
      "✦ Configure hospital profile & details",
      "✦ Auto-credential delivery to doctors"],
     ORANGE),
    ("⚙️", "Super Admin",
     ["✦ Control ALL hospitals on the platform",
      "✦ Manage labs, blood banks, & hospitals",
      "✦ View platform-wide analytics",
      "✦ Monitor patient growth trends",
      "✦ Manage collaborated partnerships"],
     RGBColor(0xFF,0x4D,0x6D)),
]

cols = [(0.4, 2.0), (3.5, 2.0), (6.6, 2.0), (9.7, 2.0)]
for (icon, role, points, color), (cx, cy) in zip(roles, cols):
    add_rect(s, cx, cy, 3.0, 5.1, fill_color=DARK_GRAY, line_color=color, line_width=1)
    add_text(s, icon, cx + 0.1, cy + 0.2, 2.8, 0.65, font_size=30, align=PP_ALIGN.CENTER)
    add_text(s, role, cx + 0.1, cy + 0.9, 2.8, 0.55,
             font_size=15, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_rect(s, cx + 0.3, cy + 1.5, 2.4, 0.04, fill_color=color)
    for j, pt in enumerate(points):
        add_text(s, pt, cx + 0.15, cy + 1.65 + j * 0.65, 2.7, 0.55,
                 font_size=10.5, color=SOFT_WHITE)

accent_bar(s, color=RGBColor(0x7B,0x61,0xFF))
slide_number(s, 7)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — TECHNOLOGY STACK
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=ORANGE)

add_text(s, "TECHNOLOGY", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=ORANGE)
add_text(s, "Built With Industry-Standard Technology", 0.5, 0.7, 12.0, 0.9,
         font_size=36, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 5.0, 0.05, fill_color=ORANGE)

tech_groups = [
    ("🖥️ Frontend", [
        ("React.js 18", "Modern UI library with component-based architecture"),
        ("Tailwind CSS", "Utility-first responsive styling framework"),
        ("React Router v6", "Client-side navigation and routing"),
        ("Framer Motion", "Smooth animations and micro-interactions"),
        ("Axios", "HTTP client for API communication"),
    ], CYAN),
    ("⚙️ Backend", [
        ("FastAPI (Python)", "High-performance async REST API framework"),
        ("MongoDB Atlas", "Scalable NoSQL cloud database"),
        ("JWT Authentication", "Secure token-based user authentication"),
        ("WebSockets", "Real-time data updates for queue system"),
        ("OpenStreetMap / Overpass", "Real-world hospital and lab discovery"),
    ], ORANGE),
    ("☁️ Infrastructure", [
        ("Render.com", "Cloud deployment with auto-scaling"),
        ("Vercel (Frontend)", "Fast CDN-based frontend delivery"),
        ("SendGrid / Brevo", "Automated email and credential delivery"),
        ("Google Maps API", "Hospital and lab navigation integration"),
        ("Twilio / SMS", "Emergency alert and notification service"),
    ], GREEN),
]

col_x = [0.4, 4.6, 8.8]
for (group, items, color), cx in zip(tech_groups, col_x):
    add_rect(s, cx, 1.9, 4.0, 5.3, fill_color=DARK_GRAY, line_color=color, line_width=1)
    add_rect(s, cx, 1.9, 4.0, 0.55, fill_color=color)
    add_text(s, group, cx + 0.15, 1.93, 3.7, 0.5,
             font_size=14, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    for j, (tech, detail) in enumerate(items):
        ty = 2.6 + j * 0.9
        add_rect(s, cx + 0.15, ty, 0.06, 0.35, fill_color=color)
        add_text(s, tech, cx + 0.35, ty, 3.5, 0.35,
                 font_size=12, bold=True, color=color)
        add_text(s, detail, cx + 0.35, ty + 0.32, 3.5, 0.45,
                 font_size=9.5, color=SOFT_WHITE)

accent_bar(s, color=ORANGE)
slide_number(s, 8)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — WHAT MAKES US UNIQUE
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)

add_text(s, "COMPETITIVE ADVANTAGE", 0.5, 0.3, 6.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "Why MediChain+ is Unique", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 4.0, 0.05, fill_color=CYAN)

# Comparison table header
headers = ["FEATURE", "Practo", "Apollo 24|7", "1mg", "MediChain+"]
col_w  = [4.2, 1.8, 1.95, 1.8, 2.0]
col_xs = [0.4, 4.65, 6.5, 8.45, 10.3]

header_y = 1.85
add_rect(s, 0.4, header_y, 12.0, 0.5, fill_color=DARK_GRAY)
for hdr, cx, cw in zip(headers, col_xs, col_w):
    clr = CYAN if hdr == "MediChain+" else SOFT_WHITE
    add_text(s, hdr, cx, header_y + 0.07, cw, 0.4,
             font_size=11, bold=True, color=clr, align=PP_ALIGN.CENTER)

rows_data = [
    ("GPS Hospital Discovery (Live)",         "❌", "❌", "❌", "✅"),
    ("Blood Bank with Blood Group Filter",     "❌", "❌", "❌", "✅"),
    ("Emergency 108 Routing + SMS Alerts",     "❌", "❌", "❌", "✅"),
    ("Dean & Super Admin Portals",             "❌", "❌", "❌", "✅"),
    ("Multi-Hospital Patient Management",      "❌", "❌", "❌", "✅"),
    ("OpenStreetMap External Facility Search", "❌", "❌", "❌", "✅"),
    ("Doctor Appointment Booking",             "✅", "✅", "✅", "✅"),
    ("AI Health Assistant",                    "❌", "✅", "❌", "✅"),
]

for i, (feat, p, ap, mg, mc) in enumerate(rows_data):
    ry = 2.4 + i * 0.55
    row_color = RGBColor(0x12,0x28,0x42) if i % 2 == 0 else DARK_GRAY
    add_rect(s, 0.4, ry, 12.0, 0.52, fill_color=row_color)
    add_text(s, feat, col_xs[0], ry + 0.08, col_w[0], 0.38,
             font_size=10.5, color=SOFT_WHITE)
    for val, cx, cw in zip([p, ap, mg, mc], col_xs[1:], col_w[1:]):
        clr = GREEN if val == "✅" else RGBColor(0xFF,0x4D,0x6D)
        if cx == col_xs[-1]:  # MediChain+ column
            add_rect(s, cx - 0.05, ry + 0.03, cw + 0.05, 0.44,
                     fill_color=RGBColor(0x00,0x3B,0x4A), line_color=CYAN, line_width=0.5)
        add_text(s, val, cx, ry + 0.08, cw, 0.38,
                 font_size=14, bold=True, color=clr, align=PP_ALIGN.CENTER)

accent_bar(s)
slide_number(s, 9)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — MARKET OPPORTUNITY
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=GREEN)

add_text(s, "MARKET OPPORTUNITY", 0.5, 0.3, 6.0, 0.5, font_size=11, bold=True, color=GREEN)
add_text(s, "A Massive & Growing Market", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 4.5, 0.05, fill_color=GREEN)

stats = [
    ("$372B", "India's healthcare\nmarket size by 2030", GREEN),
    ("600M+", "Internet users in India\nlooking for health info", CYAN),
    ("80%", "Patients prefer\nonline appointment booking", ORANGE),
    ("5X", "Growth in health-tech\nadoption post-COVID", RGBColor(0x7B,0x61,0xFF)),
]

for i, (number, label, color) in enumerate(stats):
    sx = 0.5 + i * 3.2
    add_rect(s, sx, 2.0, 2.9, 2.5, fill_color=DARK_GRAY, line_color=color, line_width=1.5)
    add_text(s, number, sx + 0.1, 2.15, 2.7, 1.1,
             font_size=42, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_text(s, label, sx + 0.1, 3.2, 2.7, 1.1,
             font_size=12, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

# Target segments
add_text(s, "🎯  Our Target Segments", 0.5, 4.7, 8.0, 0.5,
         font_size=15, bold=True, color=GREEN)

segments = [
    ("🏥 Tier-2 & Tier-3 City Hospitals", "Hospitals outside metro cities have no digital presence — a completely untapped market."),
    ("👨‍👩‍👧 Urban Families", "Families managing healthcare for multiple members need one centralized platform."),
    ("🏫 Medical Colleges", "Students and staff at medical institutions — built-in user base from day one."),
]
for i, (seg, desc) in enumerate(segments):
    sy = 5.2 + i * 0.65
    add_rect(s, 0.5, sy, 12.2, 0.58, fill_color=DARK_GRAY, line_color=GREEN, line_width=0.5)
    add_text(s, seg, 0.65, sy + 0.09, 3.5, 0.42, font_size=12, bold=True, color=GREEN)
    add_text(s, desc, 4.2, sy + 0.09, 8.3, 0.42, font_size=11, color=SOFT_WHITE)

accent_bar(s, color=GREEN)
slide_number(s, 10)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — BUSINESS MODEL
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=ORANGE)

add_text(s, "BUSINESS MODEL", 0.5, 0.3, 6.0, 0.5, font_size=11, bold=True, color=ORANGE)
add_text(s, "How MediChain+ Generates Revenue", 0.5, 0.7, 12.0, 0.9,
         font_size=36, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 5.0, 0.05, fill_color=ORANGE)

streams = [
    ("💼", "Hospital SaaS",
     "Monthly fee charged to hospitals\nfor listing and management tools.\n₹5,000 – ₹50,000 / month per hospital.",
     CYAN),
    ("🤝", "Booking Commission",
     "10–15% commission on every\ndoctor consultation booked\nthrough the platform.",
     GREEN),
    ("🔬", "Lab Partnership",
     "10–20% revenue share on\nall lab tests booked through\nMediChain+ partner labs.",
     ORANGE),
    ("⭐", "Premium Patients",
     "₹499/month patient subscription\nfor priority booking, family accounts,\nand health reminders.",
     RGBColor(0x7B,0x61,0xFF)),
]

for i, (icon, title, desc, color) in enumerate(streams):
    bx = 0.4 + i * 3.2
    add_rect(s, bx, 2.0, 3.0, 4.8, fill_color=DARK_GRAY, line_color=color, line_width=1)
    add_rect(s, bx, 2.0, 3.0, 0.1, fill_color=color)
    add_text(s, icon, bx + 0.1, 2.2, 2.8, 0.7, font_size=36, align=PP_ALIGN.CENTER)
    add_text(s, title, bx + 0.1, 3.0, 2.8, 0.55,
             font_size=14, bold=True, color=color, align=PP_ALIGN.CENTER)
    add_rect(s, bx + 0.4, 3.6, 2.2, 0.04, fill_color=color)
    add_text(s, desc, bx + 0.15, 3.75, 2.7, 2.8,
             font_size=11.5, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

accent_bar(s, color=ORANGE)
slide_number(s, 11)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — ROADMAP
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)

add_text(s, "PRODUCT ROADMAP", 0.5, 0.3, 6.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "Our Journey to Scale", 0.5, 0.7, 12.0, 0.9,
         font_size=38, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 3.5, 0.05, fill_color=CYAN)

phases = [
    ("Phase 1\nNow — Month 2", [
        "✦ 5 partner hospitals onboarded",
        "✦ Payment (Razorpay) integrated",
        "✦ Live doctor queue system launched",
        "✦ Real-time bed availability tracker",
        "✦ Progressive Web App (PWA) conversion",
    ], CYAN, "CURRENT"),
    ("Phase 2\nMonth 3–6", [
        "✦ 25 hospitals across 3 cities",
        "✦ QR Health Passport for patients",
        "✦ Vitals tracking dashboard",
        "✦ Family health manager feature",
        "✦ First 1,000 paying patients",
    ], GREEN, "NEXT"),
    ("Phase 3\nMonth 7–12", [
        "✦ Native Android & iOS applications",
        "✦ Insurance partner integrations",
        "✦ 50+ hospitals, 5 cities",
        "✦ ABDM (Govt) integration",
        "✦ Series A funding application",
    ], ORANGE, "FUTURE"),
]

for i, (phase, items, color, badge) in enumerate(phases):
    px = 0.4 + i * 4.3
    add_rect(s, px, 1.85, 4.1, 5.4, fill_color=DARK_GRAY, line_color=color, line_width=1.5)
    add_rect(s, px, 1.85, 4.1, 0.12, fill_color=color)
    add_rect(s, px + 2.8, 1.95, 1.1, 0.35, fill_color=color)
    add_text(s, badge, px + 2.8, 1.97, 1.1, 0.35,
             font_size=8, bold=True, color=NAVY, align=PP_ALIGN.CENTER)
    add_text(s, phase, px + 0.15, 2.1, 3.8, 0.7,
             font_size=15, bold=True, color=color)
    add_rect(s, px + 0.2, 2.8, 3.7, 0.04, fill_color=color)
    for j, item in enumerate(items):
        add_text(s, item, px + 0.2, 2.95 + j * 0.73, 3.7, 0.6,
                 font_size=11.5, color=SOFT_WHITE)

accent_bar(s)
slide_number(s, 12)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — VISION & IMPACT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=RGBColor(0x7B,0x61,0xFF))

add_text(s, "OUR VISION", 0.5, 0.3, 6.0, 0.5, font_size=11, bold=True,
         color=RGBColor(0x7B,0x61,0xFF))
add_text(s, "Building the Future of Healthcare in India", 0.5, 0.7, 12.0, 0.9,
         font_size=34, bold=True, color=WHITE)
add_rect(s, 0.5, 1.55, 6.5, 0.05, fill_color=RGBColor(0x7B,0x61,0xFF))

add_text(s,
         '"To become India\'s most trusted healthcare network — where every patient,\n'
         'regardless of location or economic background, has instant access\n'
         'to quality hospitals, doctors, and emergency services."',
         1.0, 1.8, 11.0, 1.5, font_size=17, italic=True, color=SOFT_WHITE,
         align=PP_ALIGN.CENTER)

impacts = [
    ("🏥\n10,000+", "Hospitals\nConnected"),
    ("👤\n5 Million+", "Patients\nServed"),
    ("🌆\n100+", "Cities\nCovered"),
    ("⏱️\n60 Seconds", "To Find &\nBook Any Doctor"),
]
for i, (num, label) in enumerate(impacts):
    ix = 0.5 + i * 3.2
    add_rect(s, ix, 3.5, 2.9, 2.2, fill_color=DARK_GRAY,
             line_color=RGBColor(0x7B,0x61,0xFF), line_width=1)
    add_text(s, num, ix + 0.1, 3.6, 2.7, 1.05,
             font_size=24, bold=True, color=RGBColor(0x7B,0x61,0xFF),
             align=PP_ALIGN.CENTER)
    add_text(s, label, ix + 0.1, 4.65, 2.7, 0.8,
             font_size=12, color=SOFT_WHITE, align=PP_ALIGN.CENTER)

add_text(s,
         "MediChain+ is not just a product — it is a mission to democratize healthcare "
         "access across India and make quality medical services available to every citizen.",
         0.5, 5.85, 12.3, 0.9, font_size=13, italic=True, color=SOFT_WHITE,
         align=PP_ALIGN.CENTER)

accent_bar(s, color=RGBColor(0x7B,0x61,0xFF))
slide_number(s, 13)


# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — THANK YOU / CONTACT
# ═══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)
add_rect(s, 0, 6.8, 13.33, 0.08, fill_color=CYAN)

add_text(s, "THANK YOU", 0.5, 0.5, 12.0, 0.55,
         font_size=13, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, "MediChain+", 0.5, 1.05, 12.0, 1.4,
         font_size=72, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, "Connecting Healthcare. Empowering Lives.", 0.5, 2.4, 12.0, 0.65,
         font_size=22, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

add_rect(s, 3.5, 3.1, 6.33, 0.05, fill_color=CYAN)

add_text(s, "Get In Touch", 0.5, 3.3, 12.0, 0.55,
         font_size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

contact_info = [
    "🌐  Website:   www.medichain.in",
    "📧  Email:     contact@medichain.in",
    "📱  WhatsApp:  +91 XXXXX XXXXX",
    "📍  Location:  Vignan University, Guntur, Andhra Pradesh",
]
add_multiline(s, contact_info, 2.5, 3.9, 8.0, 2.5,
              font_size=14, color=SOFT_WHITE, align=PP_ALIGN.CENTER, line_spacing=6)

add_text(s, "🚀  We are open for Hospital Partnerships · Investor Meetings · Pilot Programs",
         0.5, 6.15, 12.3, 0.55, font_size=13, bold=True, color=CYAN,
         align=PP_ALIGN.CENTER)

slide_number(s, 14)


# ─── Save ─────────────────────────────────────────────────────────────────────
output_path = os.path.join(os.path.dirname(__file__), "..", "MediChainPlus_Presentation.pptx")
output_path = os.path.abspath(output_path)
prs.save(output_path)
print("\n[SUCCESS] Presentation saved successfully!")
print(f"[FILE]    {output_path}")
print("[SLIDES]  14 slides generated")
print("\n  Open the file in Microsoft PowerPoint or Google Slides.")
