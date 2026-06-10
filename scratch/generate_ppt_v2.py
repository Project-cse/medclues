"""
MediChain+ — Professional PowerPoint Presentation Generator (WITH IMAGES)
Run this script to create the full PPT: python scratch/generate_ppt_v2.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# ─── Color Palette ───────────────────────────────────────────────────────────
NAVY        = RGBColor(0x0A, 0x1D, 0x37)   # Deep navy blue
CYAN        = RGBColor(0x00, 0xC2, 0xCB)   # Bright cyan
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)   # White
DARK_GRAY   = RGBColor(0x1E, 0x2D, 0x40)   # Dark gray
SOFT_WHITE  = RGBColor(0xCC, 0xE5, 0xFF)   # Soft blue-white
ORANGE      = RGBColor(0xFF, 0x8C, 0x00)   # Orange
GREEN       = RGBColor(0x00, 0xE0, 0x96)   # Mint green

IMAGE_DIR = r"C:\Users\shaik\OneDrive - Vignan University\Desktop\PMS FNL\scratch\ppt_images"

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]

def add_rect(slide, left, top, width, height, fill_color=None, line_color=None, line_width=None):
    shape = slide.shapes.add_shape(1, Inches(left), Inches(top), Inches(width), Inches(height))
    shape.line.fill.background()
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width or 1)
    else:
        shape.line.fill.background()
    return shape

def add_text(slide, text, left, top, width, height,
             font_name="Calibri", font_size=18, bold=False, italic=False,
             color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox

def add_image(slide, img_name, left, top, width=None, height=None):
    img_path = os.path.join(IMAGE_DIR, img_name)
    if os.path.exists(img_path):
        if width and height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width), height=Inches(height))
        elif width:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), width=Inches(width))
        elif height:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top), height=Inches(height))
        else:
            slide.shapes.add_picture(img_path, Inches(left), Inches(top))
    else:
        print(f"[WARNING] Image missing: {img_path}")

def bg(slide, color=NAVY):
    add_rect(slide, 0, 0, 13.33, 7.5, fill_color=color)

def accent_bar(slide, y=6.8, color=CYAN):
    add_rect(slide, 0, y, 13.33, 0.08, fill_color=color)

def slide_number(slide, num, total=14, color=SOFT_WHITE):
    add_text(slide, f"{num} / {total}", 12.5, 7.1, 0.8, 0.3,
             font_size=9, color=color, align=PP_ALIGN.RIGHT)

# ─── Slides ──────────────────────────────────────────────────────────────────

# Slide 1: Title
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_rect(s, 0, 0, 0.12, 7.5, fill_color=CYAN)
add_text(s, "MediChain+", 1.2, 1.8, 9.0, 1.6, font_size=72, bold=True, color=CYAN)
add_text(s, "India's Smartest Hospital Network Platform", 1.2, 3.2, 9.5, 0.7, font_size=26, color=WHITE)
add_image(s, "home.png", 8.0, 1.5, width=4.5) # Hero image preview
accent_bar(s)
slide_number(s, 1)

# Slide 2: The Problem
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "THE PROBLEM", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=ORANGE)
add_text(s, "Healthcare Fragmentation", 0.5, 0.7, 10.0, 0.9, font_size=38, bold=True, color=WHITE)
add_text(s, "• No central place for all health services\n• Emergency response delays\n• Blood group search is manual and slow\n• Appointment waiting times are unpredictable", 0.7, 2.5, 6.0, 4.0, font_size=22, color=SOFT_WHITE)
add_image(s, "login.png", 7.5, 2.0, width=5.0) # Show login/onboarding screen
accent_bar(s, color=ORANGE)
slide_number(s, 2)

# Slide 3: The Solution
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "THE SOLUTION", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "A Unified Healthcare Ecosystem", 0.5, 0.7, 10.0, 0.9, font_size=38, bold=True, color=WHITE)
add_image(s, "home.png", 6.5, 1.8, width=6.0) # Full home page view
add_text(s, "MediChain+ bridges the gap between\npatients and healthcare professionals,\noffering a real-time, integrated network.", 0.7, 2.5, 5.5, 3.0, font_size=20, color=SOFT_WHITE)
accent_bar(s)
slide_number(s, 3)

# Slide 4: Real-Time Discovery
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "HOSPITAL DISCOVERY", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=CYAN)
add_text(s, "Live GPS-Based Searching", 0.5, 0.7, 10.0, 0.9, font_size=38, bold=True, color=WHITE)
add_image(s, "hospitals.png", 5.5, 1.8, width=7.0) # Real hospitals list
add_text(s, "• Real-time distance calculation\n• Filter by Specialty & Type\n• OpenStreetMap Integration\n• Contact & Navigate directly", 0.7, 2.5, 4.5, 4.0, font_size=20, color=SOFT_WHITE)
accent_bar(s)
slide_number(s, 4)

# Slide 5: Doctor Appointments
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "DOCTOR BOOKING", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=GREEN)
add_text(s, "Streamlined Clinical Access", 0.5, 0.7, 10.0, 0.9, font_size=38, bold=True, color=WHITE)
add_image(s, "doctors.png", 5.5, 1.8, width=7.0) # Real doctors list
add_text(s, "• Filter by Experience & Specialty\n• Instant Slot Confirmation\n• Professional Profile Reviews\n• Seamless Patient Dashboard", 0.7, 2.5, 4.5, 4.0, font_size=20, color=SOFT_WHITE)
accent_bar(s, color=GREEN)
slide_number(s, 5)

# Slide 6: Diagnostic Labs
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "DIAGNOSTIC LABS", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=ORANGE)
add_text(s, "Laboratory & Test Management", 0.5, 0.7, 10.0, 0.9, font_size=38, bold=True, color=WHITE)
add_image(s, "labs.png", 5.5, 1.8, width=7.0) # Real labs list
add_text(s, "• Book MRI, ECG, Blood Tests\n• Filter by Ratings & Proximity\n• Collaborated Lab Network\n• Digital Test Request System", 0.7, 2.5, 4.5, 4.0, font_size=20, color=SOFT_WHITE)
accent_bar(s, color=ORANGE)
slide_number(s, 6)

# Slide 7: Admin Control Center
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "ADMINISTRATION", 0.5, 0.3, 5.0, 0.5, font_size=11, bold=True, color=RGBColor(0x7B,0x61,0xFF))
add_text(s, "Powerful Back-Office Tools", 0.5, 0.7, 10.0, 0.9, font_size=38, bold=True, color=WHITE)
add_image(s, "admin.png", 5.5, 1.8, width=7.0) # Real admin screenshot
add_text(s, "• Super Admin Oversight\n• Hospital Dean Dashboard\n• Doctor Credentialing\n• Live Booking Analytics", 0.7, 2.5, 4.5, 4.0, font_size=20, color=SOFT_WHITE)
accent_bar(s, color=RGBColor(0x7B,0x61,0xFF))
slide_number(s, 7)

# Slide 8: Thank You
s = prs.slides.add_slide(BLANK_LAYOUT)
bg(s)
add_text(s, "THANK YOU", 0.5, 3.0, 12.3, 1.0, font_size=60, bold=True, color=CYAN, align=PP_ALIGN.CENTER)
add_text(s, "www.medichain.in  |  contact@medichain.in", 0.5, 4.2, 12.3, 0.5, font_size=18, color=WHITE, align=PP_ALIGN.CENTER)
accent_bar(s)
slide_number(s, 8, total=8)

# Save
output_path = r"C:\Users\shaik\OneDrive - Vignan University\Desktop\PMS FNL\MediChainPlus_Visual_Presentation.pptx"
prs.save(output_path)
print("[SUCCESS] Visual Presentation saved!")
print(f"[FILE]    {output_path}")
